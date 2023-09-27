import streamlit as st
from PyPDF2 import PdfReader
import docx2txt
from utils.html_blocks import bot_template, user_template, css
from dotenv import load_dotenv
from pptx import Presentation
import openai
from langchain.memory import ConversationBufferMemory
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.embeddings import OpenAIEmbeddings
from langchain.vectorstores import FAISS
from langchain.chat_models import ChatOpenAI
from langchain.chains import ConversationalRetrievalChain


def document_loader(batch_files):
    text = ""
    for file in batch_files:

        if file.type == 'application/pdf':
            st.session_state['valid_file'] = True
            reader = PdfReader(file)
            for page in reader.pages:
                text += page.extract_text()

        elif file.type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
            st.session_state['valid_file'] = True
            doc =  docx2txt.process(file)
            text += doc

        elif file.type == 'text/plain':
            st.session_state['valid_file'] = True
            text += str(file.read(),"utf-8")

        elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            st.session_state['valid_file'] = True
            prs = Presentation(file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text
        else:
            st.session_state['valid_file'] = False
            st.session_state['chatbot-disabled'] = True
            return st.error('Incorrect File Type')
            
    return text

def split_text(document):
    text_splitter = RecursiveCharacterTextSplitter(
    chunk_size = 1000,
    chunk_overlap = 200,
    length_function = len)

    pages = text_splitter.split_text(document)

    return pages

def vector_storage(pages,api_key):
    embeddings = OpenAIEmbeddings(openai_api_key=api_key)
    storage = FAISS.from_texts(pages,embeddings)
    
    return storage

def get_conversation_chain(storage,api_key):
    llm = ChatOpenAI(model_name = 'gpt-3.5-turbo',openai_api_key=api_key)
    memory = ConversationBufferMemory(memory_key='chat_history', return_messages=True)
    chain = ConversationalRetrievalChain.from_llm(llm=llm, retriever=storage.as_retriever(),memory = memory)

    return chain

def handle_user_input(question):
    response = st.session_state.conversation({'question':question})
    st.session_state.chat_history = response['chat_history']
    
    for i, message in enumerate(st.session_state.chat_history):
        if i % 2 == 0:
            st.write(user_template.replace("{{MSG}}", message.content), unsafe_allow_html=True)
        else:
            st.write(bot_template.replace("{{MSG}}", message.content), unsafe_allow_html=True)